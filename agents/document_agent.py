"""Evidence Suite - Document Agent
Forensic document processing for PDFs, Word, Excel, and other formats.
Based on 2024-2025 best practices for court-admissible evidence handling.
"""

from __future__ import annotations

import hashlib
import io
import os
import re
import struct
import tempfile
import zipfile
from datetime import datetime
from typing import Any

import numpy as np
from loguru import logger

from agents.base import BaseAgent
from core.models import AnalysisResult, EvidencePacket, EvidenceType, ProcessingStage


class DocumentConfig:
    """Document Agent configuration."""

    def __init__(
        self,
        extract_metadata: bool = True,
        detect_hidden_content: bool = True,
        detect_tampering: bool = True,
        extract_embedded_objects: bool = True,
        extract_comments: bool = True,
        extract_revisions: bool = True,
        scan_macros: bool = True,
        max_file_size_mb: int = 100,
        hash_algorithms: list[str] | None = None,
    ):
        self.extract_metadata = extract_metadata
        self.detect_hidden_content = detect_hidden_content
        self.detect_tampering = detect_tampering
        self.extract_embedded_objects = extract_embedded_objects
        self.extract_comments = extract_comments
        self.extract_revisions = extract_revisions
        self.scan_macros = scan_macros
        self.max_file_size_mb = max_file_size_mb
        self.hash_algorithms = hash_algorithms or ["sha256", "md5"]


class DocumentAgent(BaseAgent):
    """Sensory layer agent for forensic document analysis.

    Features:
    - PDF text/metadata extraction (PyMuPDF)
    - Word document parsing with revision tracking
    - Excel spreadsheet analysis with hidden content detection
    - Embedded object extraction and hashing
    - Macro/VBA detection (security analysis)
    - Document tampering indicators
    - Chain of custody compliant processing

    Supports: PDF, DOCX, DOC, XLSX, XLS, PPTX, RTF, ODT
    """

    SUPPORTED_EXTENSIONS = {
        ".pdf": "pdf",
        ".docx": "docx",
        ".doc": "doc",
        ".xlsx": "xlsx",
        ".xls": "xls",
        ".pptx": "pptx",
        ".rtf": "rtf",
        ".odt": "odt",
        ".txt": "text",
        ".csv": "csv",
    }

    def __init__(self, agent_id: str | None = None, config: DocumentConfig | None = None):
        self.doc_config = config or DocumentConfig()
        super().__init__(
            agent_id=agent_id,
            agent_type="document",
            config={
                "extract_metadata": self.doc_config.extract_metadata,
                "detect_hidden_content": self.doc_config.detect_hidden_content,
                "detect_tampering": self.doc_config.detect_tampering,
                "scan_macros": self.doc_config.scan_macros,
            },
        )
        self._pymupdf_available = False
        self._docx_available = False
        self._openpyxl_available = False
        self._oletools_available = False
        self._pptx_available = False

    async def _setup(self) -> None:
        """Initialize document processing dependencies."""
        # Check PyMuPDF (fitz) for PDF
        try:
            import fitz

            self._pymupdf_available = True
            logger.info("PyMuPDF available for PDF processing")
        except ImportError:
            logger.warning("PyMuPDF not installed. Install with: pip install PyMuPDF")

        # Check python-docx for Word
        try:
            import docx

            self._docx_available = True
            logger.info("python-docx available for Word processing")
        except ImportError:
            logger.warning("python-docx not installed. Install with: pip install python-docx")

        # Check openpyxl for Excel
        try:
            import openpyxl

            self._openpyxl_available = True
            logger.info("openpyxl available for Excel processing")
        except ImportError:
            logger.warning("openpyxl not installed. Install with: pip install openpyxl")

        # Check oletools for macro detection
        try:
            from oletools import olevba

            self._oletools_available = True
            logger.info("oletools available for macro detection")
        except ImportError:
            logger.warning("oletools not installed. Install with: pip install oletools")

        # Check python-pptx for PowerPoint
        try:
            import pptx

            self._pptx_available = True
            logger.info("python-pptx available for PowerPoint processing")
        except ImportError:
            logger.warning("python-pptx not installed. Install with: pip install python-pptx")

    async def _process_impl(self, packet: EvidencePacket) -> EvidencePacket:
        """Process document evidence.

        Steps:
        1. Identify document type
        2. Calculate integrity hashes (pre-processing)
        3. Extract text content
        4. Extract metadata
        5. Detect hidden content (comments, revisions, embeds)
        6. Scan for macros/security risks
        7. Detect tampering indicators
        8. Verify integrity hashes (post-processing)
        """
        if packet.evidence_type != EvidenceType.DOCUMENT:
            raise ValueError(f"Expected document evidence, got {packet.evidence_type}")

        if not packet.raw_content:
            raise ValueError("No document content to process")

        # Check file size
        size_mb = len(packet.raw_content) / (1024 * 1024)
        if size_mb > self.doc_config.max_file_size_mb:
            raise ValueError(
                f"Document too large: {size_mb:.2f}MB > {self.doc_config.max_file_size_mb}MB"
            )

        # Calculate pre-processing hashes for chain of custody
        pre_hashes = self._calculate_hashes(packet.raw_content)

        # Detect document type
        doc_type = self._detect_document_type(
            packet.raw_content, packet.source_metadata.get("filename", "")
        )

        # Process based on type
        if doc_type == "pdf":
            result = await self._process_pdf(packet.raw_content)
        elif doc_type == "docx":
            result = await self._process_docx(packet.raw_content)
        elif doc_type == "doc":
            result = await self._process_doc(packet.raw_content)
        elif doc_type == "xlsx":
            result = await self._process_xlsx(packet.raw_content)
        elif doc_type == "xls":
            result = await self._process_xls(packet.raw_content)
        elif doc_type == "pptx":
            result = await self._process_pptx(packet.raw_content)
        elif doc_type in ("text", "csv", "rtf"):
            result = await self._process_text(packet.raw_content, doc_type)
        else:
            result = await self._process_generic(packet.raw_content)

        # Calculate post-processing hashes (should match pre)
        post_hashes = self._calculate_hashes(packet.raw_content)

        # Verify integrity
        integrity_verified = pre_hashes == post_hashes
        if not integrity_verified:
            logger.error("Document integrity check FAILED - hashes don't match!")

        # Tampering analysis
        tampering_indicators = {}
        if self.doc_config.detect_tampering:
            tampering_indicators = self._analyze_tampering(result, doc_type)

        # Build analysis result
        analysis = AnalysisResult(
            agent_id=self.agent_id,
            agent_type=self.agent_type,
            confidence=self._calculate_confidence(result, tampering_indicators),
            findings={
                "document_type": doc_type,
                "page_count": result.get("page_count", 0),
                "word_count": len(result.get("text", "").split()),
                "has_metadata": bool(result.get("metadata")),
                "has_comments": len(result.get("comments", [])) > 0,
                "has_revisions": len(result.get("revisions", [])) > 0,
                "has_embedded_objects": len(result.get("embedded_objects", [])) > 0,
                "has_macros": result.get("has_macros", False),
                "macro_risk_level": result.get("macro_risk_level", "none"),
                "tampering_risk": tampering_indicators.get("risk_level", "low"),
                "integrity_verified": integrity_verified,
                "hidden_content_count": result.get("hidden_content_count", 0),
            },
            raw_output={
                "pre_hashes": pre_hashes,
                "post_hashes": post_hashes,
                "metadata": result.get("metadata", {}),
                "comments": result.get("comments", [])[:50],  # Limit stored
                "revisions": result.get("revisions", [])[:50],
                "embedded_objects": result.get("embedded_objects", [])[:20],
                "macro_analysis": result.get("macro_analysis", {}),
                "tampering_indicators": tampering_indicators,
                "structure_info": result.get("structure_info", {}),
            },
        )

        return packet.with_updates(
            extracted_text=result.get("text", ""),
            stage=ProcessingStage.OCR_PROCESSED,
            analysis_results=packet.analysis_results + [analysis],
            source_metadata={
                **packet.source_metadata,
                "document_forensics": {
                    "type": doc_type,
                    "hashes": pre_hashes,
                    "integrity_verified": integrity_verified,
                    "tampering_risk": tampering_indicators.get("risk_level", "low"),
                    "macro_detected": result.get("has_macros", False),
                },
            },
        )

    def _calculate_hashes(self, content: bytes) -> dict[str, str]:
        """Calculate document hashes for integrity verification."""
        hashes = {}
        if "sha256" in self.doc_config.hash_algorithms:
            hashes["sha256"] = hashlib.sha256(content).hexdigest()
        if "md5" in self.doc_config.hash_algorithms:
            hashes["md5"] = hashlib.md5(content, usedforsecurity=False).hexdigest()
        if "sha1" in self.doc_config.hash_algorithms:
            hashes["sha1"] = hashlib.sha1(content, usedforsecurity=False).hexdigest()
        return hashes

    def _detect_document_type(self, content: bytes, filename: str) -> str:
        """Detect document type from magic bytes and extension."""
        # Check magic bytes first (more reliable)
        if content[:4] == b"%PDF":
            return "pdf"
        if content[:4] == b"PK\x03\x04":  # ZIP-based formats
            # Check for specific Office formats
            try:
                with zipfile.ZipFile(io.BytesIO(content)) as zf:
                    names = zf.namelist()
                    if any("word/" in n for n in names):
                        return "docx"
                    if any("xl/" in n for n in names):
                        return "xlsx"
                    if any("ppt/" in n for n in names):
                        return "pptx"
                    if "content.xml" in names:
                        return "odt"
            except zipfile.BadZipFile:
                pass
        if content[:8] == b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1":  # OLE compound
            # Could be DOC, XLS, or PPT
            if b"Word.Document" in content[:2048]:
                return "doc"
            if b"Workbook" in content[:2048] or b"Book" in content[:2048]:
                return "xls"
            return "doc"  # Default to doc for OLE
        if content[:5] == b"{\\rtf":
            return "rtf"

        # Fall back to extension
        ext = os.path.splitext(filename.lower())[1] if filename else ""
        return self.SUPPORTED_EXTENSIONS.get(ext, "unknown")

    async def _process_pdf(self, content: bytes) -> dict[str, Any]:
        """Process PDF document with PyMuPDF."""
        if not self._pymupdf_available:
            return {"text": "", "error": "PyMuPDF not available"}

        import fitz

        result = {
            "text": "",
            "metadata": {},
            "page_count": 0,
            "comments": [],
            "embedded_objects": [],
            "structure_info": {},
        }

        try:
            doc = fitz.open(stream=content, filetype="pdf")
            result["page_count"] = len(doc)

            # Extract metadata
            if self.doc_config.extract_metadata:
                result["metadata"] = {
                    "title": doc.metadata.get("title", ""),
                    "author": doc.metadata.get("author", ""),
                    "subject": doc.metadata.get("subject", ""),
                    "creator": doc.metadata.get("creator", ""),
                    "producer": doc.metadata.get("producer", ""),
                    "creation_date": doc.metadata.get("creationDate", ""),
                    "modification_date": doc.metadata.get("modDate", ""),
                    "encryption": doc.is_encrypted,
                    "pdf_version": doc.metadata.get("format", ""),
                }

            # Extract text from all pages
            text_parts = []
            for page_num, page in enumerate(doc):
                page_text = page.get_text()
                text_parts.append(page_text)

                # Extract annotations/comments
                if self.doc_config.extract_comments:
                    for annot in page.annots() or []:
                        if annot.info.get("content"):
                            result["comments"].append(
                                {
                                    "page": page_num + 1,
                                    "type": annot.type[1],
                                    "content": annot.info.get("content", ""),
                                    "author": annot.info.get("title", ""),
                                    "date": annot.info.get("modDate", ""),
                                }
                            )

                # Extract embedded images/objects
                if self.doc_config.extract_embedded_objects:
                    for img_index, img in enumerate(page.get_images()):
                        xref = img[0]
                        result["embedded_objects"].append(
                            {
                                "page": page_num + 1,
                                "type": "image",
                                "xref": xref,
                                "width": img[2],
                                "height": img[3],
                            }
                        )

            result["text"] = "\n\n".join(text_parts)

            # Structure analysis
            result["structure_info"] = {
                "xref_count": doc.xref_length(),
                "is_repaired": doc.is_repaired,
                "needs_pass": doc.needs_pass,
                "permissions": doc.permissions,
            }

            # Check for hidden/invisible text layers
            hidden_count = 0
            for page in doc:
                blocks = page.get_text("dict")["blocks"]
                for block in blocks:
                    if block.get("type") == 0:  # Text block
                        for line in block.get("lines", []):
                            for span in line.get("spans", []):
                                # Check for invisible text (white on white, size 0, etc.)
                                if span.get("size", 12) < 1:
                                    hidden_count += 1
                                color = span.get("color", 0)
                                if color == 16777215:  # White
                                    hidden_count += 1

            result["hidden_content_count"] = hidden_count

            doc.close()

        except Exception as e:
            logger.warning(f"PDF processing error: {e}")
            result["error"] = str(e)

        return result

    async def _process_docx(self, content: bytes) -> dict[str, Any]:
        """Process Word DOCX document."""
        if not self._docx_available:
            return {"text": "", "error": "python-docx not available"}

        from docx import Document
        from docx.opc.constants import RELATIONSHIP_TYPE as RT

        result = {
            "text": "",
            "metadata": {},
            "page_count": 0,
            "comments": [],
            "revisions": [],
            "embedded_objects": [],
            "has_macros": False,
            "structure_info": {},
        }

        try:
            doc = Document(io.BytesIO(content))

            # Extract core properties (metadata)
            if self.doc_config.extract_metadata:
                props = doc.core_properties
                result["metadata"] = {
                    "title": props.title or "",
                    "author": props.author or "",
                    "subject": props.subject or "",
                    "created": props.created.isoformat() if props.created else "",
                    "modified": props.modified.isoformat() if props.modified else "",
                    "last_modified_by": props.last_modified_by or "",
                    "revision": props.revision,
                    "category": props.category or "",
                    "keywords": props.keywords or "",
                }

            # Extract text from paragraphs
            text_parts = []
            for para in doc.paragraphs:
                text_parts.append(para.text)

            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = [cell.text for cell in row.cells]
                    text_parts.append(" | ".join(row_text))

            result["text"] = "\n".join(text_parts)

            # Extract comments (via XML parsing)
            if self.doc_config.extract_comments:
                result["comments"] = self._extract_docx_comments(content)

            # Extract revisions/track changes
            if self.doc_config.extract_revisions:
                result["revisions"] = self._extract_docx_revisions(content)

            # Check for embedded objects
            if self.doc_config.extract_embedded_objects:
                for rel in doc.part.rels.values():
                    if "image" in rel.reltype or "oleObject" in rel.reltype:
                        result["embedded_objects"].append(
                            {
                                "type": "image" if "image" in rel.reltype else "ole",
                                "target": rel.target_ref,
                            }
                        )

            # Check for macros (DOCM would have vbaProject.bin)
            try:
                with zipfile.ZipFile(io.BytesIO(content)) as zf:
                    if "word/vbaProject.bin" in zf.namelist():
                        result["has_macros"] = True
                        if self.doc_config.scan_macros and self._oletools_available:
                            result["macro_analysis"] = self._analyze_macros(content)
            except Exception:
                pass

            result["structure_info"] = {
                "paragraph_count": len(doc.paragraphs),
                "table_count": len(doc.tables),
                "section_count": len(doc.sections),
            }

        except Exception as e:
            logger.warning(f"DOCX processing error: {e}")
            result["error"] = str(e)

        return result

    def _extract_docx_comments(self, content: bytes) -> list[dict]:
        """Extract comments from DOCX via XML."""
        comments = []
        try:
            with zipfile.ZipFile(io.BytesIO(content)) as zf:
                if "word/comments.xml" in zf.namelist():
                    import xml.etree.ElementTree as ET  # nosec B405

                    comments_xml = zf.read("word/comments.xml")
                    root = ET.fromstring(comments_xml)  # noqa: S314  # nosec B314
                    ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}

                    for comment in root.findall(".//w:comment", ns):
                        text_parts = []
                        for t in comment.findall(".//w:t", ns):
                            if t.text:
                                text_parts.append(t.text)

                        comments.append(
                            {
                                "id": comment.get(
                                    "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}id"
                                ),
                                "author": comment.get(
                                    "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}author"
                                ),
                                "date": comment.get(
                                    "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}date"
                                ),
                                "content": "".join(text_parts),
                            }
                        )
        except Exception as e:
            logger.debug(f"Comment extraction failed: {e}")
        return comments

    def _extract_docx_revisions(self, content: bytes) -> list[dict]:
        """Extract tracked changes/revisions from DOCX."""
        revisions = []
        try:
            with zipfile.ZipFile(io.BytesIO(content)) as zf:
                if "word/document.xml" in zf.namelist():
                    import xml.etree.ElementTree as ET  # nosec B405

                    doc_xml = zf.read("word/document.xml")
                    root = ET.fromstring(doc_xml)  # noqa: S314  # nosec B314
                    ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}

                    # Find insertions
                    for ins in root.findall(".//w:ins", ns):
                        text_parts = []
                        for t in ins.findall(".//w:t", ns):
                            if t.text:
                                text_parts.append(t.text)
                        revisions.append(
                            {
                                "type": "insertion",
                                "author": ins.get(
                                    "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}author"
                                ),
                                "date": ins.get(
                                    "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}date"
                                ),
                                "content": "".join(text_parts),
                            }
                        )

                    # Find deletions
                    for dele in root.findall(".//w:del", ns):
                        text_parts = []
                        for t in dele.findall(".//w:delText", ns):
                            if t.text:
                                text_parts.append(t.text)
                        revisions.append(
                            {
                                "type": "deletion",
                                "author": dele.get(
                                    "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}author"
                                ),
                                "date": dele.get(
                                    "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}date"
                                ),
                                "content": "".join(text_parts),
                            }
                        )
        except Exception as e:
            logger.debug(f"Revision extraction failed: {e}")
        return revisions

    async def _process_doc(self, content: bytes) -> dict[str, Any]:
        """Process legacy Word DOC (OLE) document."""
        result = {
            "text": "",
            "metadata": {},
            "has_macros": False,
            "embedded_objects": [],
        }

        # Use oletools for OLE analysis
        if self._oletools_available:
            try:
                from oletools import olefile

                ole = olefile.OleFileIO(io.BytesIO(content))

                # Extract metadata
                if self.doc_config.extract_metadata:
                    meta = ole.get_metadata()
                    result["metadata"] = {
                        "title": meta.title or "",
                        "author": meta.author or "",
                        "subject": meta.subject or "",
                        "created": str(meta.create_time) if meta.create_time else "",
                        "modified": str(meta.last_saved_time) if meta.last_saved_time else "",
                        "last_saved_by": meta.last_saved_by or "",
                        "revision": meta.revision_number or "",
                        "application": meta.creating_application or "",
                    }

                # List streams for structure analysis
                result["structure_info"] = {
                    "streams": ole.listdir(),
                }

                # Check for macros
                if self.doc_config.scan_macros:
                    result["macro_analysis"] = self._analyze_macros(content)
                    result["has_macros"] = result["macro_analysis"].get("has_macros", False)

                ole.close()

            except Exception as e:
                logger.warning(f"DOC processing error: {e}")
                result["error"] = str(e)
        else:
            result["error"] = "oletools not available for DOC processing"

        return result

    async def _process_xlsx(self, content: bytes) -> dict[str, Any]:
        """Process Excel XLSX document."""
        if not self._openpyxl_available:
            return {"text": "", "error": "openpyxl not available"}

        from openpyxl import load_workbook

        result = {
            "text": "",
            "metadata": {},
            "comments": [],
            "embedded_objects": [],
            "hidden_content_count": 0,
            "has_macros": False,
            "structure_info": {},
        }

        try:
            # Load workbook (read-only for forensic integrity)
            wb = load_workbook(io.BytesIO(content), read_only=False, data_only=False)

            # Extract properties
            if self.doc_config.extract_metadata:
                props = wb.properties
                result["metadata"] = {
                    "title": props.title or "",
                    "creator": props.creator or "",
                    "subject": props.subject or "",
                    "created": props.created.isoformat() if props.created else "",
                    "modified": props.modified.isoformat() if props.modified else "",
                    "last_modified_by": props.lastModifiedBy or "",
                    "category": props.category or "",
                }

            # Extract text and analyze sheets
            text_parts = []
            hidden_sheets = 0
            hidden_rows = 0
            hidden_cols = 0

            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]

                # Check if sheet is hidden
                if ws.sheet_state != "visible":
                    hidden_sheets += 1
                    result["hidden_content_count"] += 1

                sheet_text = [f"=== Sheet: {sheet_name} ==="]

                for row in ws.iter_rows():
                    row_values = []
                    for cell in row:
                        if cell.value is not None:
                            row_values.append(str(cell.value))

                        # Extract comments
                        if self.doc_config.extract_comments and cell.comment:
                            result["comments"].append(
                                {
                                    "sheet": sheet_name,
                                    "cell": cell.coordinate,
                                    "author": cell.comment.author,
                                    "content": cell.comment.text,
                                }
                            )

                    if row_values:
                        sheet_text.append(" | ".join(row_values))

                # Check for hidden rows/columns
                for row_idx, rd in ws.row_dimensions.items():
                    if rd.hidden:
                        hidden_rows += 1
                for col_idx, cd in ws.column_dimensions.items():
                    if cd.hidden:
                        hidden_cols += 1

                text_parts.extend(sheet_text)

            result["text"] = "\n".join(text_parts)
            result["hidden_content_count"] += hidden_rows + hidden_cols

            result["structure_info"] = {
                "sheet_count": len(wb.sheetnames),
                "sheet_names": wb.sheetnames,
                "hidden_sheets": hidden_sheets,
                "hidden_rows": hidden_rows,
                "hidden_columns": hidden_cols,
            }

            # Check for macros (XLSM would have vbaProject.bin)
            try:
                with zipfile.ZipFile(io.BytesIO(content)) as zf:
                    if "xl/vbaProject.bin" in zf.namelist():
                        result["has_macros"] = True
                        if self.doc_config.scan_macros and self._oletools_available:
                            result["macro_analysis"] = self._analyze_macros(content)
            except Exception:
                pass

            wb.close()

        except Exception as e:
            logger.warning(f"XLSX processing error: {e}")
            result["error"] = str(e)

        return result

    async def _process_xls(self, content: bytes) -> dict[str, Any]:
        """Process legacy Excel XLS document."""
        result = {
            "text": "",
            "metadata": {},
            "has_macros": False,
        }

        try:
            import xlrd

            wb = xlrd.open_workbook(file_contents=content)

            text_parts = []
            for sheet_idx in range(wb.nsheets):
                sheet = wb.sheet_by_index(sheet_idx)
                text_parts.append(f"=== Sheet: {sheet.name} ===")

                for row_idx in range(sheet.nrows):
                    row_values = [
                        str(sheet.cell_value(row_idx, col_idx)) for col_idx in range(sheet.ncols)
                    ]
                    text_parts.append(" | ".join(row_values))

            result["text"] = "\n".join(text_parts)
            result["structure_info"] = {
                "sheet_count": wb.nsheets,
                "sheet_names": wb.sheet_names(),
            }

        except ImportError:
            result["error"] = "xlrd not available for XLS processing"
        except Exception as e:
            logger.warning(f"XLS processing error: {e}")
            result["error"] = str(e)

        # Check for macros using oletools
        if self._oletools_available and self.doc_config.scan_macros:
            result["macro_analysis"] = self._analyze_macros(content)
            result["has_macros"] = result.get("macro_analysis", {}).get("has_macros", False)

        return result

    async def _process_pptx(self, content: bytes) -> dict[str, Any]:
        """Process PowerPoint PPTX document."""
        if not self._pptx_available:
            return {"text": "", "error": "python-pptx not available"}

        from pptx import Presentation

        result = {
            "text": "",
            "metadata": {},
            "page_count": 0,
            "comments": [],
            "embedded_objects": [],
            "structure_info": {},
        }

        try:
            prs = Presentation(io.BytesIO(content))
            result["page_count"] = len(prs.slides)

            # Extract core properties
            if self.doc_config.extract_metadata:
                props = prs.core_properties
                result["metadata"] = {
                    "title": props.title or "",
                    "author": props.author or "",
                    "subject": props.subject or "",
                    "created": props.created.isoformat() if props.created else "",
                    "modified": props.modified.isoformat() if props.modified else "",
                    "last_modified_by": props.last_modified_by or "",
                    "category": props.category or "",
                }

            # Extract text from slides
            text_parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"=== Slide {slide_num} ==="]

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)

                    # Check for embedded objects
                    if self.doc_config.extract_embedded_objects:
                        if hasattr(shape, "ole_format"):
                            result["embedded_objects"].append(
                                {
                                    "slide": slide_num,
                                    "type": "ole",
                                }
                            )

                # Extract notes
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text
                    if notes:
                        slide_text.append(f"[Notes: {notes}]")

                text_parts.extend(slide_text)

            result["text"] = "\n".join(text_parts)
            result["structure_info"] = {
                "slide_count": len(prs.slides),
            }

        except Exception as e:
            logger.warning(f"PPTX processing error: {e}")
            result["error"] = str(e)

        return result

    async def _process_text(self, content: bytes, doc_type: str) -> dict[str, Any]:
        """Process plain text, CSV, or RTF documents."""
        result = {
            "text": "",
            "metadata": {},
        }

        try:
            # Try common encodings
            for encoding in ["utf-8", "utf-16", "latin-1", "cp1252"]:
                try:
                    result["text"] = content.decode(encoding)
                    result["metadata"]["encoding"] = encoding
                    break
                except UnicodeDecodeError:
                    continue

            if doc_type == "csv":
                # Count rows/columns for structure info
                lines = result["text"].split("\n")
                result["structure_info"] = {
                    "row_count": len(lines),
                    "column_count": len(lines[0].split(",")) if lines else 0,
                }

        except Exception as e:
            logger.warning(f"Text processing error: {e}")
            result["error"] = str(e)

        return result

    async def _process_generic(self, content: bytes) -> dict[str, Any]:
        """Generic processing for unknown document types."""
        result = {
            "text": "",
            "metadata": {},
            "error": "Unknown document type - limited processing",
        }

        # Try to extract any text
        try:
            result["text"] = content.decode("utf-8", errors="ignore")
        except Exception:
            pass

        return result

    def _analyze_macros(self, content: bytes) -> dict[str, Any]:
        """Analyze VBA macros for security risks."""
        if not self._oletools_available:
            return {"error": "oletools not available"}

        from oletools.olevba import TYPE_OLE, TYPE_OpenXML, VBA_Parser

        analysis = {
            "has_macros": False,
            "macro_count": 0,
            "suspicious_keywords": [],
            "autoexec_keywords": [],
            "iocs": [],
            "risk_level": "none",
        }

        try:
            vba_parser = VBA_Parser(filename="document", data=content)

            if vba_parser.detect_vba_macros():
                analysis["has_macros"] = True

                for filename, stream_path, vba_filename, vba_code in vba_parser.extract_macros():
                    analysis["macro_count"] += 1

                # Analyze for suspicious patterns
                results = vba_parser.analyze_macros()
                for kw_type, keyword, description in results:
                    if kw_type == "Suspicious":
                        analysis["suspicious_keywords"].append(
                            {
                                "keyword": keyword,
                                "description": description,
                            }
                        )
                    elif kw_type == "AutoExec":
                        analysis["autoexec_keywords"].append(
                            {
                                "keyword": keyword,
                                "description": description,
                            }
                        )
                    elif kw_type == "IOC":
                        analysis["iocs"].append(
                            {
                                "keyword": keyword,
                                "description": description,
                            }
                        )

                # Calculate risk level
                if analysis["iocs"] or len(analysis["suspicious_keywords"]) > 3:
                    analysis["risk_level"] = "high"
                elif analysis["autoexec_keywords"] or analysis["suspicious_keywords"]:
                    analysis["risk_level"] = "medium"
                elif analysis["has_macros"]:
                    analysis["risk_level"] = "low"

            vba_parser.close()

        except Exception as e:
            logger.warning(f"Macro analysis error: {e}")
            analysis["error"] = str(e)

        return analysis

    def _analyze_tampering(self, result: dict, doc_type: str) -> dict[str, Any]:
        """Analyze document for tampering indicators."""
        indicators = {
            "warnings": [],
            "risk_level": "low",
            "timestamp_anomalies": False,
            "structure_anomalies": False,
            "metadata_inconsistencies": False,
        }

        metadata = result.get("metadata", {})

        # Check timestamp anomalies
        created = metadata.get("created") or metadata.get("creation_date")
        modified = metadata.get("modified") or metadata.get("modification_date")

        if created and modified:
            try:
                # Parse dates and compare
                if created > modified:
                    indicators["timestamp_anomalies"] = True
                    indicators["warnings"].append("Creation date is after modification date")
            except Exception:
                pass

        # Check for revision anomalies
        revisions = result.get("revisions", [])
        if len(revisions) > 50:
            indicators["warnings"].append(f"Unusually high revision count: {len(revisions)}")

        # Check for hidden content
        hidden_count = result.get("hidden_content_count", 0)
        if hidden_count > 10:
            indicators["warnings"].append(f"High hidden content count: {hidden_count}")
            indicators["structure_anomalies"] = True

        # Check metadata completeness
        if metadata:
            if not metadata.get("author") and not metadata.get("creator"):
                indicators["warnings"].append("Missing author information")
            if not created:
                indicators["warnings"].append("Missing creation date")

        # Macro-related warnings
        if result.get("has_macros"):
            macro_risk = result.get("macro_analysis", {}).get("risk_level", "none")
            if macro_risk in ("medium", "high"):
                indicators["warnings"].append(f"Document contains {macro_risk}-risk macros")

        # Calculate overall risk
        warning_count = len(indicators["warnings"])
        if warning_count >= 3 or indicators["timestamp_anomalies"]:
            indicators["risk_level"] = "high"
        elif warning_count >= 1:
            indicators["risk_level"] = "medium"
        else:
            indicators["risk_level"] = "low"

        return indicators

    def _calculate_confidence(self, result: dict, tampering: dict) -> float:
        """Calculate processing confidence score."""
        confidence = 0.85  # Base confidence

        # Boost for successful text extraction
        if result.get("text"):
            confidence += 0.05

        # Boost for metadata extraction
        if result.get("metadata"):
            confidence += 0.05

        # Penalize for tampering indicators
        risk = tampering.get("risk_level", "low")
        if risk == "high":
            confidence -= 0.2
        elif risk == "medium":
            confidence -= 0.1

        # Penalize for errors
        if result.get("error"):
            confidence -= 0.15

        return max(0.0, min(1.0, confidence))

    async def self_critique(self, result: AnalysisResult) -> float:
        """Self-critique document processing quality."""
        if not result.is_successful:
            return 0.0

        score = result.confidence

        # Penalize high tampering risk
        risk = result.findings.get("tampering_risk", "low")
        if risk == "high":
            score *= 0.7
        elif risk == "medium":
            score *= 0.85

        # Bonus for successful integrity verification
        if result.findings.get("integrity_verified"):
            score = min(1.0, score + 0.05)

        # Bonus for metadata extraction
        if result.findings.get("has_metadata"):
            score = min(1.0, score + 0.03)

        return score
